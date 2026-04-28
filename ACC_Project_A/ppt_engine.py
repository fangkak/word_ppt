from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def md_to_ppt(md_file, ppt_file, project_info):
    """
    将 Markdown 文件转换为 PowerPoint 演示文稿
    
    Args:
        md_file: Markdown 源文件路径
        ppt_file: PowerPoint 输出文件路径
        project_info: 项目信息字典
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 创建标题幻灯片
    add_title_slide(prs, project_info)
    
    # 读取 Markdown 文件并解析
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # 按照 Markdown 结构添加幻灯片
    add_content_slides(prs, content)
    
    # 保存 PowerPoint 文件
    prs.save(ppt_file)
    print(f"PowerPoint 生成完成: {ppt_file}")

def add_title_slide(prs, project_info):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # 深蓝色背景
    
    # 项目名称
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = project_info.get("project_name", "项目名称")
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 项目描述
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_frame.text = project_info.get("description", "项目描述")
    desc_frame.paragraphs[0].font.size = Pt(18)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 项目信息
    info_text = f"客户: {project_info.get('customer', 'N/A')}\n版本: {project_info.get('version', 'N/A')}\n日期: {project_info.get('date', 'N/A')}"
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = info_text
    info_frame.paragraphs[0].font.size = Pt(12)
    info_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slides(prs, content):
    """
    解析 Markdown 内容并添加到幻灯片
    """
    lines = content.split('\n')
    current_slide = None
    
    for line in lines:
        line = line.strip()
        
        # 一级标题作为新幻灯片的标题
        if line.startswith('# '):
            title = line.replace('# ', '')
            current_slide = add_content_slide(prs, title)
        
        # 二级标题作为子标题
        elif line.startswith('## '):
            subtitle = line.replace('## ', '')
            if current_slide:
                add_subtitle_to_slide(current_slide, subtitle)
        
        # 普通文本作为内容
        elif line and not line.startswith('#') and current_slide:
            add_text_to_slide(current_slide, line)
        
        # 图片处理
        elif line.startswith('!['):
            img_path = extract_image_path(line)
            if img_path and os.path.exists(img_path):
                add_image_to_slide(current_slide, img_path)

def add_content_slide(prs, title):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    content_box.text_frame.word_wrap = True
    
    return slide

def add_subtitle_to_slide(slide, subtitle):
    """添加副标题"""
    # 获取已有的内容框
    text_frame = slide.text_frame if hasattr(slide, 'text_frame') else None
    
    # 查找现有的内容文本框
    for shape in slide.shapes:
        if shape.has_text_frame and shape.left == Inches(0.5) and shape.top == Inches(1.3):
            p = shape.text_frame.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
            p.level = 0
            break

def add_text_to_slide(slide, text):
    """添加文本内容"""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.left == Inches(0.5) and shape.top == Inches(1.3):
            p = shape.text_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(16)
            p.level = 0
            break

def add_image_to_slide(slide, image_path):
    """添加图片"""
    try:
        slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(8))
    except Exception as e:
        print(f"添加图片失败 {image_path}: {e}")

def extract_image_path(line):
    """从 Markdown 图片语法中提取路径"""
    # ![alt](path) -> path
    start = line.find('(')
    end = line.find(')')
    if start != -1 and end != -1:
        return line[start+1:end]
    return None
